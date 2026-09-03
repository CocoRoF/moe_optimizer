"""Build the results deck from the user's template masters (docs/slides/)."""
import copy, json, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
TPL = sys.argv[1] if len(sys.argv) > 1 else "/home/workspace_copy/my_paper/장하렴_템플릿.pptx"
OUT = "docs/slides/MoE_Contribution_Skipping_2026-09-03.pptx"
p = Presentation(TPL)
# keep the template's title-slide text boxes (name/affiliation) for reuse, then drop all slides
title_src = p.slides[0]
name_text = next((sh.text_frame.text for sh in title_src.shapes if sh.has_text_frame and "경희대학교" in sh.text_frame.text), "")
sldIdLst = p.slides._sldIdLst
for sldId in list(sldIdLst):
    p.part.drop_rel(sldId.rId); sldIdLst.remove(sldId)
L_TITLE, L_CONTENT, L_SECTION, L_TITLEONLY = p.slide_layouts[0], p.slide_layouts[1], p.slide_layouts[2], p.slide_layouts[5]
W, H = p.slide_width, p.slide_height
NAVY = RGBColor(0x1f, 0x4e, 0x79); RED = RGBColor(0xc0, 0x50, 0x4d); GREY = RGBColor(0x59, 0x59, 0x59)

def title_only(text):
    s = p.slides.add_slide(L_TITLEONLY); s.shapes.title.text = text
    for r in s.shapes.title.text_frame.paragraphs[0].runs: r.font.size = Pt(26)
    return s
def section(text, sub=""):
    s = p.slides.add_slide(L_SECTION); s.shapes.title.text = text
    if sub and len(s.placeholders) > 1: s.placeholders[1].text = sub
    return s
def bullets(s, items, left=0.6, top=1.35, width=12.1, height=5.6, size=16):
    height = min(height, Emu(H).inches - top - 0.35)          # never extend below the canvas
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl = 0
        while it.startswith("  "): it = it[2:]; lvl += 1
        para.level = lvl; para.text = ("• " if lvl == 0 else "– ") + it
        for r in para.runs: r.font.size = Pt(size - 2 * lvl); r.font.color.rgb = GREY if lvl else RGBColor(0x26, 0x26, 0x26)
        para.space_after = Pt(6)
    return tb
def table(s, rows, left, top, width, col_w=None, size=11, bold_rows=(), height=None):
    n, m = len(rows), len(rows[0])
    shp = s.shapes.add_table(n, m, Inches(left), Inches(top), Inches(width), Inches(height or 0.32 * n)); t = shp.table
    if col_w:
        for j, w_ in enumerate(col_w): t.columns[j].width = Inches(w_)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.text = str(val); para = c.text_frame.paragraphs[0]
            for r in para.runs:
                r.font.size = Pt(size); r.font.bold = (i == 0 or i in bold_rows)
                if i == 0: r.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
            if i == 0: c.fill.solid(); c.fill.fore_color.rgb = NAVY
            elif i in bold_rows: c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xdc, 0xe6, 0xf1)
    return shp
def picture(s, path, left, top, width=None, height=None):
    return s.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width) if width else None, height=Inches(height) if height else None)
def caption(s, text, left, top, width, size=10):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4)); tf = tb.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = text
    for r in tf.paragraphs[0].runs: r.font.size = Pt(size); r.font.italic = True; r.font.color.rgb = GREY

R = "results/"
sw = {x["policy"]: x for x in json.load(open(R + "olmoe/policy_sweep_olmoe.json"))}; base = sw["top8"]["ppl"]
d = lambda n: f"{(sw[n]['ppl']/base-1)*100:+.1f} %"
q = {x["policy"]: x for x in json.load(open(R + "qwen3/policy_sweep_qwen3.json"))}; qb = q["top8"]["ppl"]
dq = lambda n: f"{(q[n]['ppl']/qb-1)*100:+.1f} %"

# ---- 1. title
s = p.slides.add_slide(L_TITLE)
s.shapes.title.text = "Contribution-Calibrated Dynamic Expert Skipping\nfor Bandwidth-Bound MoE Decoding"
for r in s.shapes.title.text_frame.paragraphs[0].runs: r.font.size = Pt(30)
if len(s.placeholders) > 1:
    s.placeholders[1].text = "router score가 아니라 calibrated 기여도로 expert를 고른다 — training-free, router 보존, CPU 재현\n" + (name_text.split("/")[0].strip() if name_text else "") + "\n2026-09-03"
# ---- 2. 목차
s = title_only("목차")
bullets(s, ["서론 — MoE 디코드는 대역폭 제한: 실행하지 않은 expert = 읽지 않은 바이트", "이론적 배경 — 2023–2026 계보 세 갈래, training-free 스킵 규칙의 약점 W1–W4",
            "연구 설계 — 메커니즘 · 스트리밍 엔진 · matched-k′ 프로토콜 · 사전 등록 gate", "연구 결과 — OLMoE 확정(CI) · 대역폭 1.80× · tail · oracle · Qwen3 null과 그 원인",
            "기여점 · 한계 · 향후 연구", "References"], size=18)
# ---- 서론
section("서론", "Research Background & Problem")
s = title_only("MoE 추론의 병목은 연산이 아니라 메모리 대역폭이다")
bullets(s, ["fine-grained MoE(OLMoE 64/8, Qwen3-30B 128/8)는 토큰마다 8개 expert의 가중치 3개 행렬을 읽는다",
            "batch-1 디코드에서는 이 읽기가 지배적: OLMoE fp32 기준 token당 2.1 GB, Qwen3 0.86 GB (측정, F16)",
            "동적 expert 스킵 = 라우팅된 8개 중 일부만 실행 → 스킵한 expert의 바이트를 그대로 절약 (선형, 측정)",
            "기존 training-free 스킵 규칙은 전부 router score 하나로 결정한다",
            "  질문: router score는 'expert가 얼마나 기여하는가'를 알고 있는가?",
            "  답(측정): 거의 모른다 — r = +0.17 (OLMoE), −0.05 (Qwen3)"], size=17)
s = title_only("연구 목표와 제약")
bullets(s, ["목표: 학습 없이, router와 expert identity를 보존한 채, 적은 품질 손실로 expert 로드를 줄인다",
            "제약 1 — training-free: calibration 1회(forward만), 역전파·distillation 없음",
            "제약 2 — 재현: CPU만으로 전 파이프라인 실행, 모델 revision·데이터 선택·seed·환경 고정 (results/ENV.md)",
            "제약 3 — 정직한 비교: 모든 규칙을 같은 평균 k′에서, 같은 코드 경로로, paired bootstrap CI와 함께",
            "제약 4 — 사전 등록 gate: 폐기 조건 G1–G4를 실험 전에 기록"], size=17)
# ---- 이론적 배경
section("이론적 배경", "Literature 2023 → 2026 (24 papers, all verified against arXiv)")
s = title_only("계보: 세 갈래는 서로 거의 만나지 않는다")
table(s, [["갈래", "질문", "대표 (연도)", "라우팅을 바꾸는가"],
          ["어디서 가져올까", "오프로딩·캐싱·프리페치", "Mixtral-offload '23, MoE-Infinity '24, ExpertFlow '24, Fate '25, Speculating Experts '26, SpecPrefetch '26", "아니오 — 예측해서 지연을 숨김"],
          ["몇 개를 쓸까", "동적 스킵", "Lu et al. ACL'24 (k=2), LExI '25 (정적 층별 k), 2512.21911 (k>2, 중앙값), ZEDA '26 (학습)", "예 — 로드 자체를 줄임 (대역폭에 비례)"],
          ["누구와 나눌까", "batch-aware", "Opportunistic '25, SERE ICLR'26", "예 — batch ≥ 16 전제"]],
      0.5, 1.3, 12.3, col_w=[2.2, 2.2, 5.4, 2.5], size=11)
caption(s, "대역폭에 비례하는 이득은 두 번째 갈래뿐. 그 training-free 계열은 얇고(4편) 전부 score 전용.", 0.5, 4.3, 12)
s = title_only("training-free 스킵 규칙의 약점 W1–W4")
bullets(s, ["W1  router score는 dispatch 신호이지 contribution 신호가 아니다 — load-balancing 손실 아래 학습되어 균등하게 밀림",
            "W2  임계값이 품질이 아니라 분위수에 맞춰짐 — 2512.21911: β_m = 중앙값 → 각 m에서 50% 고정 스킵률, 품질 제어 없음 (m=4에서 수학 붕괴)",
            "W3  평가 무대가 이득이 나는 무대가 아니다 — speculative verification 안, batch ≥ 16 GPU; batch-1 대역폭 제한 디코드 측정 없음",
            "W4  tail 미보고 — MoEXBench('26): 평균이 domain 붕괴를 숨긴다",
            "가장 좋은 수치(ZEDA '26: expert 연산 50%↓, 1.2×)는 self-distillation 학습이 필요"], size=16)
# ---- 연구 설계
section("연구 설계", "Mechanism · Engine · Protocol · Gates")
s = title_only("메커니즘: score × calibrated output scale")
bullets(s, ["calibration 1회: 층 l, expert e에 대해  s[l,e] = E‖E_e(x)‖  (토큰이 e로 라우팅될 때) — 층당 E개 float",
            "추론: top-8 후보를  c_e = w_e · s[l,e]  로 정렬, 누적 기여 비율 ≥ 1 − τ_l 인 최소 접두사만 실행",
            "τ_l: 층별 이분탐색(40회)으로 calibration 토큰에서 평균 실행 수 = 목표 k′",
            "공정 대조군 score-only = 같은 코드, s ≡ 1  →  차이는 순전히 정렬 신호",
            "baseline: 정적 top-k′ / 발표된 중앙값 규칙(2512.21911) / oracle(진짜 per-token ‖E_e(x)‖, 배포 불가 진단)",
            "학습 없음 · router 보존 · expert identity 보존 · 저장량 E floats/층"], size=16)
s = title_only("스트리밍 CPU 엔진과 실험 프로토콜")
table(s, [["항목", "값"],
          ["엔진", "layer-streaming, safetensors mmap에서 층 단위 읽기, fp32; 참조(bf16 HF) 대비 top-1 일치 100 % (OLMoE NLL 3.207 vs 3.212; Qwen3 3.508 vs 3.503)"],
          ["모델", "OLMoE-1B-7B @6d84c48 (norm_topk_prob=False) · Qwen3-30B-A3B @ad44e77 (norm_topk_prob=True)"],
          ["calibration", "WikiText-2 train 첫 2,048 (OLMoE) / 4,096 (Qwen3) token"],
          ["평가", "WikiText-2 test 첫 N token, 512-token 시퀀스; GSM8K·HumanEval 1,024 token (tail)"],
          ["통계", "시퀀스 단위 paired bootstrap, B=5000, seed 0, 95 % 구간"],
          ["디코드", "batch 1, KV 캐시, prefill 32 + 64 step; 바이트 = 엔진이 읽은 fp32 바이트"],
          ["자원", "CPU 11 threads, GPU 없음, 엔진 상주 3.7–4.5 GB; 전 파이프라인 scripts/reproduce.sh"]],
      0.5, 1.3, 12.3, col_w=[1.6, 10.7], size=11)
s = title_only("사전 등록 gate와 결과")
table(s, [["gate", "폐기 조건", "결과"],
          ["G1", "층 내 s의 CV < 0.15 (스케일이 균일 → 신호 없음)", "OLMoE 0.255, Qwen3 0.341 — 통과"],
          ["G2", "matched k′에서 contribution ≤ score-only", "OLMoE 통과 (CI가 0 제외) · Qwen3 실패 (동률)"],
          ["G3", "bytes/token 절감이 tok/s로 안 이어짐", "통과 (1.80×)"],
          ["G4", "worst-domain 저하 > 평균의 3배", "통과 (tail/mean ≤ 1.10)"]],
      0.5, 1.3, 12.3, col_w=[0.8, 6.5, 5.0], size=12)
# ---- 연구 결과
section("연구 결과", "OLMoE confirmed with CIs · bandwidth · tail · oracle · Qwen3 null")
s = title_only("결과 1 — OLMoE-1B-7B, matched k′, 8,192 token (F20)")
picture(s, "docs/figures/fig1_ppl_vs_k_olmoe.png", 0.4, 1.25, height=5.2)
table(s, [["policy", "k′≈5", "k′≈4"], ["정적 top-k", d("top5(static)"), d("top4(static)")], ["score-only (공정)", d("score_only@k'=5.0"), d("score_only@k'=4.0")],
          ["contribution (ours)", d("contribution@k'=5.0"), d("contribution@k'=4.0")], ["중앙값 규칙 (발표)", "+309 % (k′ 3.2)", "—"]],
      7.4, 1.4, 5.5, col_w=[2.5, 1.5, 1.5], size=11, bold_rows=(3,))
bullets(s, ["paired bootstrap (16 seq):", "  contribution vs score-only  −3.0 % [−4.0, −2.0]  ·  −7.0 % [−8.8, −5.3]", "  contribution vs static  −0.8 % n.s.  ·  −2.9 % [−4.9, −1.1]",
            "score-only 동적 규칙은 정적 절단보다 나쁘다 — 이득 전부가 calibrated scale에서 나온다"], left=7.4, top=3.4, width=5.6, height=3, size=12)
s = title_only("결과 2 — batch-1 디코드: 바이트는 k′에 선형, 속도 1.80× (F16)")
picture(s, "docs/figures/fig3_decode_bandwidth.png", 0.4, 1.25, height=5.0)
bullets(s, ["token당 바이트 = 16층 × k′ × 3행렬 × 2048×1024 fp32 (+ attention)", "  expert 하나 스킵 = 16.8 MB/token, 정확히 선형", "KV 캐시 경로: uncached logits와 오차 0.000",
            "contribution@5: 1521 MB/tok, 3.35 tok/s = 1.80× (top-8 1.86 tok/s)", "이 표의 63-token ppl 열은 품질 측정이 아님 — 품질은 결과 1"], left=7.6, top=1.4, width=5.4, height=5, size=13)
s = title_only("결과 3 — 왜 score로는 안 되는가 (F14, F18)")
picture(s, "docs/figures/fig2_score_vs_scale.png", 0.4, 1.25, width=8.2)
bullets(s, ["층 내 출력 스케일 CV: 0.255 (OLMoE), 0.341 (Qwen3)", "r(s_e, gate weight): +0.17 / −0.05", "  Qwen3는 48층 중 15층에서 음(−0.35까지)", "두 base 모델 모두 top-8 내 라우팅이 평탄 (head/tail ≈ 3×)",
            "score만으로 정렬하면 '잡음'(OLMoE) 또는 '역방향'(Qwen3)으로 정렬한다"], left=8.8, top=1.4, width=4.3, height=5, size=12)
s = title_only("결과 4 — tail: 수학·코드가 텍스트보다 덜 저하 (F19)")
table(s, [["policy", "wikitext", "gsm8k", "code", "tail / wikitext"], ["score-only @5", "+8.4 %", "+6.4 %", "+6.9 %", "1.00"], ["contribution @5", "+5.4 %", "+5.3 %", "+6.0 %", "1.10"],
          ["score-only @4", "+24.8 %", "+16.1 %", "+13.5 %", "1.00"], ["contribution @4", "+15.9 %", "+12.5 %", "+11.3 %", "1.00"]],
      0.5, 1.4, 12.3, col_w=[3.0, 2.3, 2.3, 2.3, 2.4], size=13, bold_rows=(2, 4))
bullets(s, ["G4 통과: worst-domain / mean ≤ 1.10 — MoEXBench가 경고한 평균 뒤의 붕괴 없음", "모든 셀에서 contribution 승", "wikitext 열은 새 1,024-token slice — F20 결과의 독립 재현 (−2.8 %, −7.1 %)"], top=3.6, size=15)
s = title_only("결과 5 — oracle: proxy가 한계인가, 신호가 한계인가 (F24)")
picture(s, "docs/figures/fig4_oracle.png", 0.4, 1.25, width=8.3)
bullets(s, ["oracle = 모든 expert를 계산하고 진짜 w_e‖E_e(x)‖로 선택 (배포 불가, 상한)", "OLMoE: oracle vs proxy +0.1 % [−1.1, +1.4] — 64-float 캘리브레이션이 달성 가능한 이득 전부를 회수",
            "Qwen3: oracle vs score-only +5.1 % [+1.7, +9.9] — 신호 자체가 실패", "→ 두 체제를 가르는 변수: router의 top-k renormalisation"], left=8.9, top=1.4, width=4.3, height=5, size=12)
s = title_only("결과 6 — Qwen3-30B-A3B: 특성화된 null (F21, F22)")
table(s, [["policy (k′=5, 4,096 token)", "Δ ppl vs top-8 11.879"], ["정적 top-5", dq("top5(static)")], ["score-only", dq("score_only@k'=5.0")], ["contribution", dq("contribution@k'=5.0")],
          ["contribution_renorm (오차모델)", dq("contribution_renorm@k'=5.0")], ["중앙값 규칙", "×15"]], 0.5, 1.3, 6.0, col_w=[4.0, 2.0], size=12)
bullets(s, ["contribution vs score-only: +0.4 % [−1.7, +2.3] — 동률 (1,024-token 캘리브레이션에서는 +5.3 % 손실 → 4배로 해소: calibration starvation)",
            "budget hogging 기각(층별 k′ 4.9–5.2 평탄) · renorm 오차모델 오히려 악화 · oracle도 실패",
            "해석: norm_topk_prob=True → expert 하나를 빼면 생존자 전부가 W_all/W_P로 재조정. 출력 변화는 방향에 의존하고, expert 출력은 근사 직교일 뿐(whitened NN cos 0.40)",
            "OLMoE(renorm 없음)에서는 제거 = 뺄셈 → norm이 충분한 proxy", "반사실 F25 (renorm off): −40 % — 그러나 모델이 ppl 42.9(4×)로 손상된 상태의 아티팩트", "**F25b (full-mass renorm; top-8 동일 11.879, 제거 = 뺄셈): +0.8 % [−0.9, +2.7] 동률, static이 여전히 우세 → renorm 가설 기각**", "가설 4개 중 3개 측정으로 기각, 1개(calibration 부족)는 F21 손실의 원인으로 확인 — null은 실재하고 기제는 미해명"], left=6.8, top=1.3, width=6.2, height=5.5, size=12)
# ---- 기여·한계·향후
section("연구의 기여점 및 향후 연구 방향")
s = title_only("기여")
bullets(s, ["측정: router score는 expert 출력 크기를 담지 않는다 — 두 fine-grained MoE에서 r ≈ 0, Qwen3는 층의 1/3에서 음", "방법: 층당 E floats의 calibrated scale로 정렬 — 학습 없음, router 보존",
            "결과(OLMoE): score-only 대비 −3.0 % / −7.0 % (CI 0 제외), 정적 top-k 대비 k′≈4에서 −2.9 %, oracle과 동률, tail 저하 ≤ 평균, 디코드 1.80×",
            "경계: Qwen3에서는 어떤 norm 기반 규칙도(proxy·oracle·renorm 중립화) score를 못 이김 — 경계는 측정됨, 기제는 미해명 (가설 3개 기각)", "부정 결과 명시: 발표된 중앙값 규칙은 두 모델 모두 붕괴; 제가 유도한 원리적 변형 둘은 휴리스틱보다 열등",
            "재현: CPU 전용, 모델 revision·데이터·seed·환경 고정, results/에 원자료, reproduce.sh 한 줄"], size=15)
s = title_only("한계와 향후 연구")
bullets(s, ["절대 비용: 로드 38 % 절감에 ppl +10 % — ZEDA(학습)의 50 % 거의 무손실과 격차 존재", "positive 모델 1개 + null 모델 1개; 일반성은 'router 유형' 주장이며 반사실(F25) 대기",
            "perplexity만 — GPU 없이는 다운스트림 정확도 불가", "batch-1 fp32 CPU 무대 — GPU 배치 서빙에서는 bytes→latency가 sublinear",
            "향후: (1) 방향/중복을 보는 skip 기준 — Qwen3류에서 score가 norm을 이기는 이유의 후보  (2) instruct/추론 모델의 'certain head' 프로파일에서 재측정  (3) GPU 다운스트림 평가  (4) 세 번째 모델(DeepSeek-V2-Lite)로 경계 확인"], size=15)
section("감사합니다")
s = title_only("References")
refs = ["Lu, X. et al. (2024). Not All Experts are Equal: Efficient Expert Pruning and Skipping for MoE LLMs. ACL 2024.", "Chitty-Venkata, K. T. et al. (2025). LExI: Layer-Adaptive Active Experts for Efficient MoE Inference. arXiv:2509.02753.",
        "Oncescu, C.-A. et al. (2025). Opportunistic Expert Activation: Batch-Aware Expert Routing for Faster Decode Without Retraining. arXiv:2511.02237.", "(2025). Accelerate Speculative Decoding with Sparse Computation in Verification. arXiv:2512.21911.",
        "Chen, Y. et al. (2026). Certain Head, Uncertain Tail: Expert-Sample for Test-Time Scaling in Fine-Grained MoE. arXiv:2602.02443.", "Lv, X. et al. (2026). Post-Trained MoE Can Skip Half Experts via Self-Distillation (ZEDA). arXiv:2605.18643.",
        "Wu, J. et al. (2026). SERE: Similarity-based Expert Re-routing for Efficient Batch Decoding in MoE Models. ICLR 2026. arXiv:2602.07616.", "Zhong, S. et al. (2024). AdapMoE: Adaptive Sensitivity-based Expert Gating and Management. arXiv:2408.10284.",
        "Madan, V. et al. (2026). Speculating Experts Accelerates Inference for Mixture-of-Experts. arXiv:2603.19289.", "Kong, J. et al. (2026). SpecPrefetch: Parameter-Efficient Expert Prefetching for Sparse MoE. arXiv:2607.24787.",
        "Benazir, A. et al. (2026). Benchmarking Composable Compression Techniques in MoE LLMs (MoEXBench). arXiv:2608.21693.", "Liu, J. et al. (2024). A Survey on Inference Optimization Techniques for Mixture of Experts Models. arXiv:2412.14219.",
        "Muennighoff, N. et al. (2024). OLMoE: Open Mixture-of-Experts Language Models. arXiv:2409.02060.  ·  Qwen Team (2025). Qwen3 Technical Report. arXiv:2505.09388."]
bullets(s, refs, size=11)
p.save(OUT); print("saved", OUT, "slides:", len(p.slides))
